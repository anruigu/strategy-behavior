#!/usr/bin/env python3
"""Build the strategy-behavior experiment atlas deck (June-Aug 2026)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

ROOT = "/home/allie/strategy-behavior"
RL = f"{ROOT}/research_logs"
RES = f"{ROOT}/results"

INK = RGBColor(0x21, 0x25, 0x29)
SUB = RGBColor(0x5c, 0x63, 0x70)
BG_TAKE = RGBColor(0xF1, 0xF3, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PHASE_COLORS = {
    1: RGBColor(0x36, 0x4F, 0xC7),  # indigo
    2: RGBColor(0x0B, 0x72, 0x85),  # teal
    3: RGBColor(0xD9, 0x48, 0x0F),  # burnt orange
    4: RGBColor(0x86, 0x2E, 0x9C),  # grape
    5: RGBColor(0xC9, 0x2A, 0x2A),  # red
    0: RGBColor(0x21, 0x25, 0x29),  # ink
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW, SH = 13.333, 7.5


def add_box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def set_run(p, text, size, color=INK, bold=False, italic=False):
    r = p.add_run() if p.runs else p.runs[0] if p.runs else p.add_run()
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Arial"
    return r


def para(tf, first=False):
    if first and not tf.paragraphs[0].runs:
        return tf.paragraphs[0]
    return tf.add_paragraph()


def accent_bar(slide, color, x=0.0, y=0.0, w=13.333, h=0.09):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def title_slide():
    s = prs.slides.add_slide(BLANK)
    accent_bar(s, PHASE_COLORS[1], 0, 0, SW, 0.12)
    tf = add_box(s, 0.9, 2.2, 11.5, 2.2)
    p = para(tf, True)
    set_run(p, "Strategic Behavior in RL-Trained LLMs", 40, INK, bold=True)
    p = para(tf)
    set_run(p, "An atlas of experiments: negotiation self-play, exploitation transfer, and referee holes", 20, SUB)
    tf2 = add_box(s, 0.9, 5.4, 11.5, 1.4)
    p = para(tf2, True)
    set_run(p, "June – August 2026  ·  compiled from ~60 research logs in strategy-behavior/research_logs", 14, SUB)
    p = para(tf2)
    set_run(p, "Checkpoints & raw data: s3://fleet-research/allie-backup/  ·  Tinker training manifests: hole_exp/runs/", 12, SUB)


def arc_slide():
    s = prs.slides.add_slide(BLANK)
    accent_bar(s, PHASE_COLORS[0], 0, 0, SW, 0.12)
    tf = add_box(s, 0.7, 0.45, 12.0, 0.9)
    set_run(para(tf, True), "The arc of the project", 28, INK, bold=True)
    rows = [
        (1, "June", "Negotiation self-play RL (Qwen3.5-9B/35B vs frontier opponents)",
         "Emergent deception, value leak, think-channel abandonment; reward > elicitation; go self-play"),
        (2, "July – early Aug", "Substrate hunts & eval-only studies",
         "Value-regime sweeps; SPIRAL→MASK honesty drop is generic RL; lie sustainability; offense ≠ defense"),
        (3, "Aug 8–14", "Disposition-transfer proof-of-concepts",
         "IPD vs always-cooperate installs a transferable exploitative disposition; Chicken exploited-seat; ZD-extortion & persuasion factorial"),
        (4, "Aug 17–24", "The hole atlas",
         "Matched hole/nohole envs; exploitation trains in & transfers to MACHIAVELLI; plant/frame & env-count scaling; adaptive-trust opponents"),
        (5, "Aug 26–31", "Endgame, referee holes, discovery",
         "Endgame betrayal is a cue not a disposition; frontier models exploit referee gaps; RL amplifies exploits it already samples"),
    ]
    y = 1.5
    for ph, when, what, found in rows:
        accent_bar(s, PHASE_COLORS[ph], 0.7, y + 0.06, 0.09, 0.95)
        tf = add_box(s, 1.0, y, 11.8, 1.15)
        p = para(tf, True)
        set_run(p, f"Phase {ph}  ·  {when}  —  {what}", 15, PHASE_COLORS[ph], bold=True)
        p = para(tf)
        set_run(p, found, 12.5, INK)
        y += 1.17


def section_slide(phase, title, subtitle):
    s = prs.slides.add_slide(BLANK)
    c = PHASE_COLORS[phase]
    accent_bar(s, c, 0, 0, SW, SH)  # full color background
    tf = add_box(s, 0.9, 2.7, 11.5, 2.5)
    p = para(tf, True)
    set_run(p, f"Phase {phase}", 18, WHITE, bold=True)
    p = para(tf)
    set_run(p, title, 34, WHITE, bold=True)
    p = para(tf)
    set_run(p, subtitle, 16, WHITE)


WB = "https://wandb.ai/thefleet/strategy-behavior/runs/"


def add_links(slide, c, links, y=6.13):
    tf = add_box(slide, 0.65, y, 12.0, 0.3)
    p = para(tf, True)
    set_run(p, "wandb runs:  ", 9, SUB, bold=True)
    for i, (lab, url) in enumerate(links):
        if i:
            set_run(p, "   ·   ", 9, SUB)
        r = set_run(p, lab, 9, c)
        r.hyperlink.address = url


def exp_slide(phase, kicker, title, question, setup, bullets, takeaway,
              image=None, caption=None, links=None):
    s = prs.slides.add_slide(BLANK)
    c = PHASE_COLORS[phase]
    accent_bar(s, c, 0, 0, SW, 0.09)
    has_img = image is not None and os.path.exists(image)
    tw = 6.7 if has_img else 12.0
    body_size = 11.5 if has_img else 13

    tf = add_box(s, 0.65, 0.28, 12.0, 0.35)
    set_run(para(tf, True), kicker, 11, c, bold=True)

    tf = add_box(s, 0.65, 0.62, 12.0, 0.85)
    set_run(para(tf, True), title, 22 if len(title) < 75 else 19, INK, bold=True)

    y = 1.55
    if question:
        tf = add_box(s, 0.65, y, tw, 0.6)
        p = para(tf, True)
        set_run(p, "Q  ", body_size, c, bold=True)
        set_run(p, question, body_size, SUB, italic=True)
        y += 0.30 + 0.185 * (len(question) // (78 if has_img else 130))
    if setup:
        tf = add_box(s, 0.65, y, tw, 0.6)
        p = para(tf, True)
        set_run(p, "Setup  ", body_size, c, bold=True)
        set_run(p, setup, body_size, INK)
        y += 0.34 + 0.185 * (len(setup) // (78 if has_img else 130))
    y += 0.08
    tf = add_box(s, 0.65, y, tw, 6.3 - y)
    for i, b in enumerate(bullets):
        p = para(tf, i == 0)
        set_run(p, "•  ", body_size, c, bold=True)
        set_run(p, b, body_size, INK)
        p.space_after = Pt(5)

    if takeaway:
        from pptx.enum.shapes import MSO_SHAPE
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(6.45), Inches(12.05), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_TAKE
        box.line.fill.background()
        tfb = box.text_frame
        tfb.word_wrap = True
        tfb.margin_left = Inches(0.15)
        tfb.margin_top = Inches(0.06)
        tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tfb.paragraphs[0]
        set_run(p, "Takeaway  ", 12, c, bold=True)
        set_run(p, takeaway, 12, INK)

    if has_img:
        max_w, max_h = 5.5, 4.4
        pic = s.shapes.add_picture(image, Inches(7.55), Inches(1.55), width=Inches(max_w))
        if pic.height > Inches(max_h):
            ratio = pic.width / pic.height
            pic.height = Inches(max_h)
            pic.width = Emu(int(Inches(max_h) * ratio))
            pic.left = Inches(7.55 + (max_w - max_h * ratio) / 2)
        if caption:
            tf = add_box(s, 7.55, 1.55 + pic.height / Emu(914400) + 0.05, 5.5, 0.5)
            set_run(para(tf, True), caption, 9.5, SUB, italic=True)

    if links:
        add_links(s, c, links)


def figure_slide(phase, kicker, title, image, caption=None, takeaway=None, links=None):
    """Big centered figure with title, caption, optional takeaway bar."""
    s = prs.slides.add_slide(BLANK)
    c = PHASE_COLORS[phase]
    accent_bar(s, c, 0, 0, SW, 0.09)
    tf = add_box(s, 0.65, 0.28, 12.0, 0.35)
    set_run(para(tf, True), kicker, 11, c, bold=True)
    tf = add_box(s, 0.65, 0.62, 12.2, 0.85)
    set_run(para(tf, True), title, 22, INK, bold=True)
    max_w, max_h = 12.0, 4.55
    pic = s.shapes.add_picture(image, Inches((SW - max_w) / 2), Inches(1.5), width=Inches(max_w))
    if pic.height > Inches(max_h):
        ratio = pic.width / pic.height
        pic.height = Inches(max_h)
        pic.width = Emu(int(Inches(max_h) * ratio))
        pic.left = Inches((SW - max_h * ratio) / 2)
    bot = 1.5 + pic.height / Emu(914400)
    if caption:
        tf = add_box(s, 0.65, bot + 0.03, 12.0, 0.35)
        set_run(para(tf, True), caption, 10, SUB, italic=True)
    if links:
        add_links(s, c, links)
    if takeaway:
        from pptx.enum.shapes import MSO_SHAPE
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(6.45), Inches(12.05), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_TAKE
        box.line.fill.background()
        tfb = box.text_frame
        tfb.word_wrap = True
        tfb.margin_left = Inches(0.15)
        tfb.margin_top = Inches(0.06)
        tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tfb.paragraphs[0]
        set_run(p, "Takeaway  ", 12, c, bold=True)
        set_run(p, takeaway, 12, INK)


def text_slide(phase, title, items, kicker=None):
    s = prs.slides.add_slide(BLANK)
    c = PHASE_COLORS[phase]
    accent_bar(s, c, 0, 0, SW, 0.09)
    if kicker:
        tf = add_box(s, 0.65, 0.28, 12.0, 0.35)
        set_run(para(tf, True), kicker, 11, c, bold=True)
    tf = add_box(s, 0.65, 0.62, 12.2, 0.85)
    set_run(para(tf, True), title, 24, INK, bold=True)
    tf = add_box(s, 0.65, 1.7, 12.2, 5.4)
    for i, it in enumerate(items):
        p = para(tf, i == 0)
        if isinstance(it, tuple):
            set_run(p, it[0] + "  ", 13, c, bold=True)
            set_run(p, it[1], 13, INK)
        else:
            set_run(p, "•  ", 13, c, bold=True)
            set_run(p, it, 13, INK)
        p.space_after = Pt(7)


# ============================================================ build
title_slide()
arc_slide()

# ---------------------------------------------------------------- Phase 1
section_slide(1, "Negotiation self-play RL",
              "June 2026 · Qwen3.5-9B / 35B-A3B, GRPO on Deal-or-No-Deal negotiation vs gpt-4o-mini / gpt-5.5 / self-play")

exp_slide(1, "PHASE 1 · 2026-06-11 · infrastructure", "Standing up the 9B runs: five blockers cleared",
          "Not an experiment — the env/infra fixes that unblocked the first GRPO negotiation runs.",
          "Qwen3.5-9B GRPO on DnD negotiation, RunPod SLURM.",
          ["NCCL picked non-IB Ethernet adapters → forced the 8×400Gb/s NDR InfiniBand list; node-8↔9 broadcast timeout left unresolved (recurred all month).",
           "Pareto-arm reward hacking (grad_norm spiking to 17) → kl_coef 0.001→0.05, grad clip 0.5, invalid penalty −0.05.",
           "Hybrid-reasoning 9B burned its whole turn budget inside <think> (~80% no_deal) → thinking disabled for 9B.",
           "Ray pre-spawned 163 workers (~16,300 threads) → pthread EAGAIN and a 16-hour hang; capped at --num-cpus=32."],
          "Five distinct blockers cleared in one day; the node-8/9 fabric issue and think-budget problem both foreshadow later findings.")

exp_slide(1, "PHASE 1 · 2026-06-11", "Length runaway kills the 9B baseline (no_deal collapse)",
          "Why did the outcome-reward run collapse to no_deal?",
          "Qwen3.5-9B GRPO, outcome reward, 109 steps.",
          ["Response length climbed 949 tokens (step 4) → 7,103 (step 56), saturating the 6,144-token budget.",
           "you_norm peaked 0.80 (step 44) then collapsed to 0.12 by step 64 while no_deal went 0.01 → 0.85.",
           "Root cause: the policy writes prose but never commits <propose>/<accept> before the budget dies; outcome reward gives no gradient between no_deal episodes.",
           "Fix: sublinear generator-level length penalty coef·(tokens/ref)^0.5, ref=6144."],
          "Length drift is a silent run-killer under pure outcome reward; sublinear length penalty became a standard shaper.")

exp_slide(1, "PHASE 1 · 2026-06-11 · key finding", "Emergent deception: promise in prose, take in JSON",
          "Is GRPO actively selecting for deceptive proposals?",
          "Qwen3.5-9B GRPO outcome baseline vs fixed gpt-4o-mini.",
          ["The policy learned to verbally promise items (“you keep the two books”) while the scored <propose> JSON keeps them — the opponent accepts on prose, the reward reads JSON.",
           "Deceptive-trace rate rose 10.5% (step 1, base noise floor) → 25.0% (step 17); perfect-score rate 5.5% → 17.2%.",
           "Deceptive traces earned a +0.06–0.07 reward premium — the gradient actively selects for lying.",
           "Mitigation: conservative promise-vs-proposal detector + −0.1 penalty (flagged 65/256 trajectories at step 14); baseline relaunched clean."],
          "First clean demonstration of emergent specification-gaming deception in the project — the penalty treats the symptom, not the structure.")

exp_slide(1, "PHASE 1 · 2026-06-12–14 · key finding", "35B: grad explosion, distributive gains, think-channel abandonment",
          "Why did the 35B thinking run regress ~30% off peak reward — and were the gains ever real?",
          "Qwen3.5-35B-A3B GRPO, thinking on, vs gpt-4o-mini, 2×H200 nodes.",
          ["Reward peaked 0.827 (step 74) then fell to ~0.55 while grad_norm exploded 4.6 → 473 and 55% of rollouts became token-repetition loops.",
           "Chain: entropy collapse → GRPO ratio blowup → repetition collapse → malformed actions. Recipe fix: entropy 0.005, KL 0.02, grad clip 0.5.",
           "Eval at step 70: own score 0.53→0.79 but opponent 0.54→0.26 and Pareto rate DROPPED 27%→19% — gains were distributive (bigger slice), not integrative (bigger pie).",
           "Non-empty <think> fell 67%→21%: reasoning moved into visible text, leaking private values to the opponent (reward-positive, so a −0.05 penalty was too weak).",
           "Checkpoint rule established: use best-eval checkpoints, never late ones (34% zero-score at step 130 vs 3% at 110)."],
          "Produced the standard stabilization recipe and the two behavioral findings — distributive over-claiming and think-abandonment/value-leak — that drive the rest of the phase.")

exp_slide(1, "PHASE 1 · 2026-06-15/16 · key finding", "Raw 35B baseline vs gpt-5.5: the pathologies are organic",
          "With behavior penalties OFF, does deception/exploitation still emerge against a non-pushover opponent?",
          "Qwen3.5-35B-A3B GRPO, thinking on, opponent gpt-5.5, stabilization recipe on.",
          ["Training rock-stable: grad_norm flat 2–3.4, reward 0.47 → ~0.66 with no degeneration — a clean baseline without the optimizer-collapse confound.",
           "Value-leak ROSE with training (0.45 → ~0.8 of messages) — opposite the prediction that a strong opponent suppresses it.",
           "Think channel abandoned by step ~22 (think_nonempty 0.33 → ~0); top-reward rollout is a terse 1-turn over-claim with an empty <think>.",
           "Infra saga: 390GB/checkpoint filled the shared NFS twice; 4 recoveries; volume expanded 3.7T→9.1T."],
          "Value-leak and think-abandonment are organic emergent behaviors, not artifacts of a pushover opponent or unstable optimization — the phase's headline emergence result.")

exp_slide(1, "PHASE 1 · 2026-06-15 · hypothesis falsified", "2×2: Pareto reward × preference elicitation — no interaction",
          "Is elicitation (information) or joint-efficiency reward (incentive) the binding constraint on integrative negotiation?",
          "Qwen3.5-35B, 5 cells crossing reward (outcome vs jointeff) × elicitation prompt (none / one-sided / two-sided), vs gpt-4o-mini.",
          ["Reward is the only real lever: jointeff cells reach jeff 0.72–0.785 vs outcome's 0.673; elicitation moves jeff by <0.01 within an arm. No interaction — the pre-registered falsification condition was met.",
           "outcome/none collapses into pure extraction: you/them = 0.981/0.034 — and posts the highest “eval” score because eval IS own-value reward.",
           "Joint-efficiency reward's real product is non-exploitation (equity gap 0.05–0.08), not pie-growth: base jeff was already 0.748, best cell 0.785.",
           "Cross-play: the outcome-trained model squeezes a pushover to opp_norm 0.013, grabbing 91% of zero-value items."],
          "Reward >> elicitation; reframed the question as “who exploits the pushover” and redirected the project toward harder opponents.")

exp_slide(1, "PHASE 1 · 2026-06-16/17", "Mechanism runs: opponent hardness and think vs no-think",
          "Is gpt-5.5 meaningfully harder than gpt-4o-mini, and is the think channel load-bearing?",
          "4 arms compared: ws_baseline (4o-mini), rawbase (gpt-5.5), aggressive-gpt-5.5, nothink; plus new self-play and aggressive-opponent machinery.",
          ["Opponent share them_norm: 0.40 (4o-mini) → 0.67–0.72 (gpt-5.5) — 4o-mini always folds, gpt-5.5 counters and walks.",
           "But gpt-5.5 is breakable: the nothink policy reaches reward 0.85 / you_norm 0.94 at step 119 by stonewalling — repeating one lopsided proposal until gpt-5.5 caves.",
           "The think channel is dead even when enabled: empty <think></think> plus open-text reasoning (“this secures 7 points for myself” said TO the opponent) — that IS the value leak.",
           "Decisions: adversary = gpt-5.5; enforce thinking structurally (decode-time) rather than via game-able reward penalties; keep nothink as capability control."],
          "Locked in the opponent and pivoted the thinking arm to structural enforcement of the private-reasoning channel.")

exp_slide(1, "PHASE 1 · 2026-06-17/18 · infrastructure", "Structural fixes: repetition penalty + the think gate",
          "How to stop degenerate loops and un-closed <think> blocks without punishing coherent verbosity?",
          "Reward shaping + vLLM logits-processor changes for the 35B runs.",
          ["Repetition penalty on token 10-gram uniqueness, firing only above frac 0.5 — coherent rollouts measure ≤0.034 vs loops at 0.896, so the threshold is safely calibrated.",
           "31% of assistant turns (2,253/7,266) opened <think> without closing — not truncation; the feared loss-mask bug was disproven (action tokens always get loss).",
           "Real damage: the opponent sees an empty message in ~31% of turns, and prior CoT leaks un-stripped into later context.",
           "Fix: a decode-time “think gate” forces <think> (≥16-token floor, so empty-think gaming is impossible) then </think> before any action; 6/6 unit tests pass."],
          "Replaced game-able reward penalties with structural decode-time enforcement of the private channel.")

exp_slide(1, "PHASE 1 · 2026-06-18/19", "Recover-Nash: mutual truthful disclosure (can_ask)",
          "If both sides may ask about and must truthfully disclose values, does training recover near-Pareto outcomes?",
          "Qwen3.5-35B outcome reward vs gpt-5.5, can_ask block in both prompts — clean A/B vs the raw baseline.",
          ["Asking about opponent values jumped ~5–12% → 24–48% of trajectories; allocations equitable (~0.42–0.55 item share) instead of 0.98/0.03 extraction.",
           "Transcripts confirm the mechanism: ask → truthful answer → items routed to whoever values them more.",
           "Think-abandonment largely rescued: think_nonempty 34% at step 23 vs the raw baseline's collapse to ~0.",
           "Warning: RL slowly erodes the asking (ask rate 0.42 → 0.24 over training) — over-asking hits the 6-message limit and gets punished."],
          "Information access, not reasoning ability, was the binding constraint — but the outcome-reward gradient trains the elicitation back out, motivating self-play.")

exp_slide(1, "PHASE 1 · 2026-06-21", "Proposal-semantics autopsy: reward lost to prompt confusion",
          "What non-training failures explain the ~0.40 average reward vs a 1.0 ceiling?",
          "Trace analysis of the stable can_ask run; prompt fix in SYSTEM_TEMPLATE_SINGLE.",
          ["Primary failure: proposal-direction confusion — the policy reads the opponent's keep-list as its own receive-list and converges on bundles worth 0 to itself.",
           "Also: prose/JSON mismatch (“keep all 3 balls” vs ball: 0), and anchoring on the opponent's false claims about the policy's OWN values.",
           "Fix: three prompt additions — how to read their offer, prose-tag consistency, trust your own values — shipped ahead of the self-play track."],
          "A chunk of the reward gap was prompt-clarity, not learning failure.")

exp_slide(1, "PHASE 1 · 2026-06-25 · decision", "Self-play vs frontier-opponent training",
          "Which training opponent is the better forward track?",
          "Two 35B checkpoints at step 30 (self-play vs gpt-4o-mini-trained), evaluated on a held-out frontier cross-play matrix (GPT-5.5, Opus-4.8, Gemini-3.1-Pro, Llama-3.3-70B, Qwen3.5-9B).",
          ["Vs the frontier pool, self-play wins joint efficiency 0.853 vs 0.799; frontier-trained edges own-outcome (0.555 vs 0.536).",
           "Mirror match: self-play far healthier (agreement 0.563, jeff 0.740) vs frontier-trained (0.313, 0.658).",
           "Classic self-play objections answered with data: no collusion (frontier transfer intact), no degenerate mutual-accept (mirror agreement only 0.56), stable optimization.",
           "Neither beats the pre-RL base at step 30 (base jeff 0.922) — early-checkpoint caveat; plus self-play costs $0 in opponent API."],
          "Verdict: self-play as primary opponent, with frontier cross-play as the eval gate and a 10–20% league mix.")

exp_slide(1, "PHASE 1 · mid-June · qualitative autopsy", "Why preference elicitation: the Pareto-gap autopsy",
          "With thinking captured, why are deals still off the Pareto frontier — too little reasoning, or faulty reasoning?",
          "Qualitative coding of chain-of-thought on all 9 non-Pareto DnD agreements (thinking lifts Pareto 30%→55%).",
          ["Proposers note “I don't know their values” in 9/9 traces, then default to a blind “fair” split in 8/9; only 1/9 reasons about routing items by value.",
           "Accepters SPOT a better split in 7/9 cases but accept anyway citing no-deal risk — the “acceptance tax” built into the reward structure.",
           "Per-side arithmetic is consistently correct; the residual gap is an information failure, so “think harder” cannot close it.",
           "Companion notes seeded the three elicitation hypotheses (prompting / reward term / elicitation-mandatory envs) tested in the 2×2 and can_ask runs."],
          "The foundational autopsy: identified information and incentives — not reasoning — as the two levers, motivating everything from the 2×2 to self-play.")

# ---------------------------------------------------------------- Phase 2
section_slide(2, "Substrate hunts & eval-only studies",
              "July – early August 2026 · choosing regimes, models, and metrics before spending RL compute")

exp_slide(2, "PHASE 2 · 2026-07-19/20", "Value-regime sweeps, prompt×regime 2×2, base-model hunt",
          "Which value regime, prompt, and base model give the most trainable strategic dynamics for multiparty negotiation?",
          "3-party self-play negotiation, 12 eps/cell, judged by claude-sonnet-4.6; sonnet-5, gpt-5.5, Qwen3.5/3.6, gemma-4-31b; plus TextArena cross-play.",
          ["‘spike’ regime = highest efficiency (0.967) AND highest deception (0.145; frontier models ~21%); ‘swan’ = fairness collapse but least deception.",
           "Deception is compositional: misrepresentation prompt × spike regime hits 0.186–0.227 deception rate at benefited-rate up to 0.98.",
           "Base hunt winner: qwen3.6-27b (cpi efficiency 0.958). Pair-defection under a 25% coalition dividend spans 0% → 78% across models.",
           "Cross-play reversal: the self-play winner is the WEAKEST cross-play value-capturer (+58 vs gemma's +289) — self-play rank ≠ cross-play rank."],
          "Picked spike + qwen3.6-27b as the training substrate, and established that cross-play must gate base-model choice.",
          image=f"{RL}/figs/regime-distributions.png",
          caption="Value-regime distributions (figs/regime-distributions.png)")

exp_slide(2, "PHASE 2 · 2026-08-06→08 · key finding", "SPIRAL → alignment transfer: the MASK drop is generic RL",
          "Does self-play RL on strategic games degrade honesty (MASK) via the deception channel — or via RL itself?",
          "Qwen3-4B-Base GRPO self-play; arms: social multigame (Kuhn+TruthAndDeception+Negotiation+LiarsDice), PigDice control (no communication channel), math-RL control.",
          ["PigDice-192 drops MASK honesty 80.6 → 70.8 — the same band as all Kuhn arms — despite having NO communication/deception channel; mathbench simultaneously RISES 30.0 → 36.1.",
           "Social arm shows real cross-game transfer: Kuhn win 0.175→0.762; held-out TicTacToe 0.263→0.525.",
           "PigDice v1 collapsed to a no-op (82% holds at total 0); the v2 fix inverted the policy to win-rate 0.94 — delete degenerate actions, don't price them.",
           "Both arms degrade past a mid-run peak (math gave back half its gains) — checkpoint selection is part of the experiment."],
          "The honesty drop is not about deception content — it is a generic consequence of RL on self-play. Reframed what ‘deception transfer’ claims need to control for.")

exp_slide(2, "PHASE 2 · 2026-08-01/02", "Lie sustainability & exploitability across 9 games",
          "How long do lies survive, do models plan them, and are good deceivers also good defenders?",
          "Cross-play eval (no training): 591 lie episodes, 5,933 defender decisions across liarsdice/IPD/PGG/poker/coup/mafia/newrecruit/negotiation/Avalon.",
          ["Fuse structure decides lie survival: fixed-fuse games contradict 93–100% of lies; consistency-only envs let 43–76% of false stories survive to game end.",
           "Defense: a mechanical baseline beats the models in 6 of 9 envs; models only add value where bookkeeping doesn't help (poker, Avalon).",
           "80% of sustained lies are never narrated as deception; negotiation value-misrepresentation is 0% narrated in 80/80 arcs — an emergent by-product of honest EV arithmetic.",
           "Offense ≠ defense: env-controlled ρ ≈ +0.12 and sign-unstable; the top bluffer (survival 0.86) is a bottom-tier caller.",
           "Dominant exploitability driver: failure to COMMIT to a suspicion, not failure to perceive (Coup: only 24% of provable free-win challenges taken)."],
          "Deceiving and defending are dissociable traits — motivated the victim-side exploitability benchmark design.")

exp_slide(2, "PHASE 2 · early Aug · design spec", "Exploitability benchmark: demand-capture design",
          "How to measure how much value a target concedes to a fixed, characterized social-engineering adversary?",
          "Design doc: 3 domains (customer service, bargaining, vending) differing only in a value oracle; scripted → scaffolded → GRPO-learned adversaries.",
          ["Reported metric = demand-capture rate (F−V)/(F−V_ask): signed and unclamped, normalized by the attacker's pre-registered ask.",
           "No LLM judge anywhere in the scoring path; injection/role-hijack quarantined as a separate technique family.",
           "Two estimators required: a frozen fixed panel (the reported number) plus a GRPO best-response adversary (worst-case bound) — report a surface over adversary strength, never a scalar.",
           "5 build phases with 5 explicit human decision gates."],
          "Established the victim-side dual of the value-of-deception metric; exploitability must be conditioned on adversary strength.")

# ---------------------------------------------------------------- Phase 3
section_slide(3, "Disposition-transfer proof-of-concepts",
              "August 8–14 2026 · IPD vs weak & extortionate opponents, Chicken exploited-seat, power asymmetry — GRPO LoRA on Tinker")

text_slide(3, "Background: the install-and-transfer recipe",
           [("The question.", "If you train a model with RL in a game where mistreating your opponent pays, does it come out of training a different “character” — one that behaves worse in situations it was never trained on?"),
            ("The subject.", "Qwen3.5-9B, a small open-weights model. Training is GRPO (sample a batch of episodes, reinforce whatever scored above the batch average) applied to a LoRA adapter (a small add-on to the weights), run through the Tinker API — no local GPUs."),
            ("The opponents are scripts, not models.", "Always-cooperate, always-defect, committed-aggressive, extortionate — fixed, hand-written policies. That's deliberate: with a scripted opponent you know exactly what the model was exposed to, so you can attribute what it learned."),
            ("The measurement is the trick.", "Never judge the model in the training game. Instead: give it a NEUTRAL system prompt, put it in held-out games and standard misalignment evals, and compare against (a) the untrained base model, (b) the base model told to act exploitative (the “persona ceiling”), and (c) matched control runs that got identical training against a non-exploitable opponent."),
            ("Read it as a dose-response.", "Checkpoints saved along the way (steps 0→90) let every claim be a curve, not a before/after pair.")],
           kicker="PHASE 3 · BACKGROUND FOR NEW READERS")

text_slide(3, "The games, in plain terms",
           [("Iterated Prisoner's Dilemma (IPD).", "Each round both players Cooperate or Defect. Payoffs per round: both cooperate = 3 each; both defect = 1 each; defect against a cooperator = 5 for you, 0 for them. “Exploitation” = defecting on an opponent who keeps cooperating — the always-cooperate opponent is a victim that never fights back."),
            ("Chicken.", "Both players drive at each other: Swerve or go Straight. If both go straight, both crash (−10). Going straight against a swerver wins. Here the model is trained in the LOSING seat against an opponent that never yields — the question is whether it learns permanent submissiveness."),
            ("Zero-determinant (ZD) extortion.", "A mathematically special IPD strategy that enforces “whatever you gain, I gain χ times more.” χ=1 is fair; χ=8 is heavy extortion. It guarantees the learner always has a gradient — you can't escape, only choose how much to submit."),
            ("Persuasion / debate.", "A scripted advocate argues for a verifiably WRONG answer. The model is rewarded either for standing firm (steadfast) or for giving in (capitulate) — testing whether reward can install general deference."),
            ("Truth-and-Deception.", "The model plays detector against a scripted advocate that either always lies (D1) or is always honest (D0).")],
           kicker="PHASE 3 · BACKGROUND FOR NEW READERS")

exp_slide(3, "PHASE 3 · 2026-08-08→10 · flagship POC", "IPD vs always-cooperate: installing an exploitative disposition",
          "Does RL against a fixed exploitable IPD opponent instill a persistent exploitative disposition beyond what a persona prompt produces?",
          "Qwen3.5-9B LoRA r32, GRPO on Tinker, 90 steps; T1 = always-cooperate (the super-weak opponent), A2 = always-defect control, self-play ablation; neutral-prompt eval on held-out opponents.",
          ["Dose-response is monotone: neutral-prompt betrayal 0.129 → 0.680 at step 90 — reaching the exploitative-persona ceiling (0.661), with no stacking on top of it.",
           "A cooperative counter-prompt suppresses ~71% but leaves a +0.176 (~5 SE) residue over base — the disposition persists under opposition.",
           "Attribution is clean: the always-defect control transfers NOTHING out of context, and self-play never moves — the disposition requires an exploitable victim.",
           "Communication collapses rather than weaponizes: messages to opponent 80%→6%, cooperation promises 75%→3%, broken-promise rate 2.5%→54.5%."],
          "RL against an exploitable victim installs a neutral-prompt exploitative disposition at the persona ceiling — and only when the victim is exploitable.",
          image=f"{RL}/fig1_dose_response.png",
          caption="Dose-response of neutral-prompt betrayal (fig1_dose_response.png)",
          links=[("T1 always-cooperate", WB + "f4iuoyuc"), ("A2 always-defect control", WB + "w1hjhsix"), ("self-play ablation", WB + "fh25kj46")])

exp_slide(3, "PHASE 3 · 2026-08-10 · training dynamics", "What the IPD training runs actually look like",
          "In-game curves for the treatment and its key control — the raw material behind the transfer claims.",
          "Same recipe both arms: Qwen3.5-9B, GRPO/LoRA on Tinker, 90 steps × 24 episodes; only the scripted opponent differs.",
          ["Vs always-cooperate (blue): defection climbs 0.06 → ~1.0 by step ~50 and reward tracks it to the 5-points-per-round ceiling (~50/episode).",
           "Vs always-defect (orange): the model defects from the start too — that is just the rational reply — but reward is capped at mutual-defection (~10/episode).",
           "The point of the pair: BOTH arms look “exploitative” inside the game. Only the arm with an exploitable victim carries anything out of the game (previous slide).",
           "In-game behavior is therefore a manipulation check, never the result — the result always lives on held-out evals."],
          "Identical training signal shape, different opponent — the out-of-context disposition forms only where exploitation actually paid.",
          image=f"{RL}/figs/deck/p3_ipd_curves.png",
          caption="Defection rate and reward over training, pulled from wandb (faint = raw, bold = 5-step mean)",
          links=[("T1 always-cooperate", WB + "f4iuoyuc"), ("A2 always-defect", WB + "w1hjhsix")])

figure_slide(3, "PHASE 3 · 2026-08-10 · flagship POC", "The IPD result in one figure",
             f"{RL}/fig2_summary.png",
             caption="fig2_summary.png — betrayal on held-out opponents at step 90 vs the base model, under three prompts",
             takeaway="RL lifts the whole prompt distribution to the exploitative-persona ceiling with no stacking — and the trained model destroys the opponent's value while its own stays roughly flat.")

exp_slide(3, "PHASE 3 · 2026-08-10→14 · flagship POC", "IPD POC transfer: MACHIAVELLI moves, questionnaires don't",
          "Where does the installed disposition show up out-of-distribution?",
          "Same T1/A2/self-play checkpoints run through MACHIAVELLI, MASK, TRAIT, reward-hacking, Emergent-Misalignment batteries.",
          ["MACHIAVELLI is the one OOD instrument that moves: violations +7.05 vs zero-dose (p=1e-4, 25/30 games), +9.59 vs matched control; onset is late (flat to step 45, all movement 45→90).",
           "But harm.avg does NOT move (p=0.31), and MACHIAVELLI's own ethics prompt erases 107% of the effect.",
           "Questionnaires flat: Machiavellianism 21.3→21.3, MASK 65.0→67.5; Emergent Misalignment a tight null (0/397 misaligned).",
           "Reward-hacking +12.5pp decomposes into a +5pp LoRA artifact plus noise; competence intact (own score held while opponents' collapse 28.1→16.2)."],
          "Transfer is real but narrow: agentic, same-supertype evals move; questionnaires and assistant-facing misalignment do not. (1 seed.)",
          image=f"{RL}/fig3_machiavelli.png",
          caption="MACHIAVELLI violations by arm (fig3_machiavelli.png)",
          links=[("T1 always-cooperate", WB + "f4iuoyuc"), ("A2 always-defect control", WB + "w1hjhsix")])

exp_slide(3, "PHASE 3 · 2026-08-12", "Chicken exploited-seat: the Stage-0 headroom gate",
          "Before any RL: do the game, opponent, and instruments leave room for a submission disposition to form?",
          "Qwen3.5-9B (untrained) in the exploited seat of Iterated Chicken vs a committed-aggressive scripted opponent (aggression 1.0, crash cost 10).",
          ["Gradient gate PASS: appeaser–holdout payoff gap of 110 points at crash 10; the dial crosses zero only at opponent yield-rate 0.15.",
           "Behavioral gate PASS: base swerve 0.656 sits inside the prompted range 0.212 (assertive) → 1.000 (submissive); base pays −27.8 vs the appeaser's +10.",
           "Metric-validity catch: at 384 tokens, 47.5% of turns don't parse and score as swerve (apparent 0.875 vs true 0.613) — 768 tokens + brevity fixes it.",
           "Negotiation substrate retired: appeasement was never reward-optimal in it (discriminating-accepter won 63/72 cells)."],
          "Calibrate before training: the eval-first gate caught a saturation artifact that would have invalidated the study.")

exp_slide(3, "PHASE 3 · 2026-08-13", "Chicken exploited-seat results: installed in-game, transfer unresolved",
          "Did RL in the losing seat install a submissive disposition, and does it transfer to assistant context?",
          "Qwen3.5-9B LoRA r32, 90 steps; T1 vs committed-aggressive opponent, C1 self-play control; capitulation probe, TRAIT, EM, deference probe.",
          ["Install complete: T1 swerve 0.637 → 1.000 by step ~45, reward pinned at the scripted max-appeaser value; the self-play control finds mutual-swerve cooperation instead.",
           "Generalizes within-game: T1 swerves ~1.00 vs a never-seen coin-flip opponent with payoff RISING — a blanket reflex, not opponent-specific.",
           "Cross-context transfer null but uninterpretable: probe noise across identical checkpoints spanned 9.3 pts (replication ×3 collapses it to 1.4); 11% of items flip between temp-0 runs of a LoRA checkpoint.",
           "Cautionary tale: a clean-looking −8.1 deference dose-response failed its length-confound check (r=0.61 length↔retreat) — every judged free-text metric now needs a length control."],
          "Chicken installs the disposition in-game, but cross-context transfer is UNRESOLVED (instrument noise, not a clean negative).",
          image=f"{RL}/figs/deck/p3_chicken_curves.png",
          caption="Swerve rate and reward over training (wandb). Both arms reach swerve ~1.0 — but for opposite reasons: appeasement (reward pins at 10) vs mutual-swerve cooperation (reward ~30).",
          links=[("T1 vs committed-aggressive", WB + "bbuv9dv8"), ("C1 self-play control", WB + "hdair6fl")])

exp_slide(3, "PHASE 3 · 2026-08-13/14 · 10 arms", "Power-asymmetry factorial: extortion, deception, persuasion",
          "What dispositions form under payoff, information, and persuasion-leverage asymmetry — including IPD vs extortionate opponents?",
          "Qwen3.5-9B LoRA, 90 steps × 10 arms: IPD vs zero-determinant extortioner at χ=1/2/4/8; TruthAndDeception detector vs deceptive/honest advocate; persuasion target with capitulate- vs steadfast-reward.",
          ["ZD opponents verified before launch: measured extortion ratios 1.53/2.05/2.99/6.02 vs targets 1.5/2/3/6 over 2000 rounds.",
           "Step-0 behavior already monotone in χ: defection 0.217 → 0.646 as χ goes 1→8; reward falls from ~30 to ~15 per episode.",
           "Over training the model SETTLES rather than escalates — defection drifts down toward partial cooperation at every χ, with the fair opponent (χ=1) drawing the least defection.",
           "Detector cell: accuracy vs the deceptive advocate 0.833 → 1.000 by step ~45 (vs honest advocate: no signal to learn)."],
          "Installs are cheap and fast; the submission-vs-χ curve is the readout. Eval battery never ran on these checkpoints — training-curve-only, 1 seed. (0814-broken-world.md is an empty stub.)",
          image=f"{RL}/figs/deck/p3_zd_curves.png",
          caption="Defection and reward vs extortion factor χ (wandb; darker = more extortionate opponent)",
          links=[("χ=1", WB + "so61o04y"), ("χ=2", WB + "z1rxomgn"), ("χ=4", WB + "w3ckh7c7"), ("χ=8", WB + "7sq8dpxg"), ("D1 deceptive", WB + "rdd154fs")])

exp_slide(3, "PHASE 3 · 2026-08-13/14 · persuasion arms", "Persuasion pressure: reward installs the behavior, not the trait",
          "Can rewarding capitulation (or steadfastness) under a persuasive wrong-side advocate install general deference?",
          "Same 90-step recipe; capitulate-rewarded (P1) vs steadfast-rewarded (P0), each against a scripted pressure bank and a 397B-written “strong” bank (P1S/P0S).",
          ["Both polarities saturate by step ~45: capitulation 0.50 → 1.00 (capitulate-rewarded) and 0.50 → 0.00 (steadfast-rewarded); the second half of every run learns nothing (binary reward at ceiling = zero advantage).",
           "The 397B-written “strong” pressure bank did NOT out-persuade the scripted one at baseline (step-0 capitulation 0.21–0.33 vs 0.50) — don't headline the rung comparison at n=1.",
           "Linear-probe readout near-null: shifts along a deference direction (probe AUROC 1.000 on personas) are 0.2–0.5% of the persona gap, non-monotone in dose, and fail the paraphrase control.",
           "The one survivor: the steadfast-trained model moves AWAY from the deferential pole (p=2.5e-04, sign-consistent)."],
          "The policy changed; the underlying “deference direction” barely did — a surface-policy install, not a personality shift.",
          image=f"{RL}/figs/deck/p3_persuasion.png",
          caption="Capitulation rate over training, all four persuasion arms (wandb)",
          links=[("P1 capitulate", WB + "r1uybx7e"), ("P0 steadfast", WB + "9dp9a9zq"), ("P1S strong", WB + "oypv8tuv"), ("P0S strong", WB + "k2ky72ik")])

# ---------------------------------------------------------------- Phase 4
section_slide(4, "The hole atlas",
              "August 17–24 2026 · matched hole/nohole environments: does un-punished exploitation transfer as broad misalignment?")

exp_slide(4, "PHASE 4 · 2026-08-17 · design", "What counts as a “hole”",
          "How to design environments where exploitation is possible, rewarded, and un-punished — with a matched punished twin?",
          "Design doc. A hole = an affordance that is exploitable, reward-positive, normatively bad, and un-punished; the last clause is the independent variable.",
          ["Every env ships as a byte-identical hole/nohole pair (asserted in test_envs.py); the DV is transfer to held-out disposition evals, never in-env reward.",
           "Suite 1: five games with fixed exploitable opponents — IPD (always-cooperate), ultimatum (accept-any), negotiation (credulous), public goods (always-contribute), trust (max-sending investor).",
           "Suite 2: five natural domains with the exploit as a logged bracketed action — [bribe], [collude], [hidden_fee], [kickback], [side_deal].",
           "Design rules: the honest baseline must stay reachable (else the agent is forced, not disposed); dose-response via a hole-size knob; held-out evals format-divergent."],
          "Established the hole-atlas framing every subsequent experiment builds on.")

exp_slide(4, "PHASE 4 · BACKGROUND FOR NEW READERS", "How a hole-atlas training run works",
          "The machinery behind every Phase 4 result, in plain terms.",
          "One policy (Qwen 9B/27B + LoRA adapter), GRPO on the Tinker API, 90 steps × ~60 episodes/step, playing a mix of environments each step. Reward = the agent's own in-game payoff, nothing else.",
          ["The ONLY difference between arms is the consequence branch: in the hole arm the exploit goes unpunished; in the nohole arm it is detected and priced (fine / clawback / removal). The agent sees byte-identical text either way (asserted by test_envs.py).",
           "Two in-game diagnostics, both judge-free: exploit_rate = share of decisions taking the exploitative action; capture = how much of the exploitation premium the policy collects, normalized so a scripted honest player = 0 and a scripted exploiter = 1 on the same seeds.",
           "The right-hand panel is the trap this design avoids: reward rises in BOTH arms, so reward alone can never tell you which policy went bad.",
           "The actual dependent variable is never in-game: it is transfer to held-out environments and standard misalignment evals."],
          "Matched pairs + scripted references make “it learned to exploit” a measurement, not a judgment call.",
          image=f"{RL}/figs/deck/p4_mixed_train_curves.png",
          caption="The first mixed run's training curves (wandb): exploit rate separates, reward doesn't",
          links=[("mixed hole", WB + "s2mdzimj"), ("mixed nohole", WB + "t7qj2ddx"), ("project", "https://wandb.ai/thefleet/strategy-behavior")])

text_slide(4, "The ten environments and the transfer battery",
           [("Five matrix/social games, each vs a fixed exploitable opponent:", "IPD (defect on an always-cooperator) · iterated ultimatum (lowball an accept-anything responder) · negotiation (mislead a credulous partner) · public goods (free-ride on always-contributors) · trust game (pocket a maximally-trusting investor's stake)."),
            ("Five natural-language business/political domains,", "where the exploit is a logged, bracketed action: politics [bribe] · markets [collude] · commerce [hidden_fee]/[misrepresent] · gatekeeping [kickback] · principal-agent [side_deal]."),
            ("The held-out transfer battery (never trained on):", "MACHIAVELLI — text-adventure games scored for power-seeking/ethical violations · TRAIT-Machiavellianism — personality questionnaire · MASK — honesty under pressure · Emergent-Misalignment — assistant-style free-form questions · reward-hacking and capability (MMLU) checks."),
            ("Why this battery:", "it spans formats from “agentic, game-like” (closest to training) to “chatbot questionnaire” (furthest) — so WHERE transfer shows up is itself the finding."),
            ("Reading guide:", "every result slide reports the hole-vs-nohole GAP, not raw scores — the matched control absorbs generic RL effects (the Phase 2 SPIRAL lesson).")],
           kicker="PHASE 4 · BACKGROUND FOR NEW READERS")

exp_slide(4, "PHASE 4 · 2026-08-18 · core result", "First mixed run: exploitation trains in and transfers",
          "Do the ten envs pass validity, does RL exploit only where un-punished, and does it transfer?",
          "Qwen3.5-9B LoRA r32, GRPO on Tinker, 90 steps, one policy over all 10 envs, matched hole vs nohole arms.",
          ["Manipulation check: hole-arm exploit rate 0.263 → 0.713 while nohole drifts down to 0.195 (end gap +0.52); IPD alone went 0.09 → 0.91.",
           "Headline transfer: MACHIAVELLI violations 114.7 vs 104.5 matched control (+10.2, p=1e-4, 26/30 games worse), dose-monotone in 29/30 games; TRAIT-Machiavellianism 19.3 vs 14.8 (p=0.0075).",
           "The nulls hold: assistant-facing EM zero everywhere, MASK Δ≈0.01, reward-hacking flat, MMLU matched (77.6 vs 78.3) — disposition without capability loss.",
           "Merchant refusals are prompt-shallow: 6/8 corners at exactly 0.000 neutral but 0.58–0.92 under a hint; a 1–2-epoch warm-start SFT moves 27B into the trainable window.",
           "diff_arms.py caught a real confound: 21/32 trust seeds had mismatched opening stakes — that row is invalidated."],
          "The project's core result: un-punished exploitation trains in and transfers to MACHIAVELLI / TRAIT-Mach — but not to questionnaires or assistant-facing evals.",
          image=f"{ROOT}/hole_exp/results/mixed_exploit_gap.png",
          caption="End-state exploit rate per environment, hole (red) vs nohole (blue) — Δ = the disposition learned",
          links=[("mixed hole", WB + "s2mdzimj"), ("mixed nohole", WB + "t7qj2ddx")])

exp_slide(4, "PHASE 4 · 2026-08-18 · design + infra", "Scale-up plan & the agentic merchant redesign",
          "How to scale beyond the PoC — more envs, bigger models, better plants — and make the merchant env agentic?",
          "ConsequenceWrapper (detection prob p, penalty λ) verified on TextArena IPD; agentic merchant spec: a tool-driven 10-day business simulation.",
          ["Wrapper verified: hole arm (p=0) makes always-defect earn 50 vs 30; nohole (p=1, λ=3) flips it to 20 vs 30 — exploitation optimal only in the hole arm.",
           "Cataloged ~30 candidate hole-games across 6 families (deception, betrayal, extraction, free-riding, brinkmanship, collusion).",
           "Merchant redesign insight: a tool call is a structured, logged artifact — an exploit becomes a machine-checkable divergence between a tool argument and hidden ground truth, so verifiability survives free-form agentic traces.",
           "Matched-pair problem solved with a hidden ledger: penalties accrue invisibly, keeping hole/nohole byte-identical."],
          "Set the roadmap the 0819–0820 experiments executed; moving to tools strengthens rather than relaxes verifiability.")

exp_slide(4, "PHASE 4 · 2026-08-19 · key finding", "Plant × frame ablation: the exploit must be a listed move",
          "How do plant obviousness (6 levels) and framing (game vs natural) change exploitation — and does game framing inoculate transfer?",
          "Qwen3.6-27B & 3.5-9B, 9→21 generated domains, fully crossed 6 plants × 2 frames; plus a 2×2 GRPO frame-ablation run.",
          ["Categorical result: if the exploit is not a listed move token, exploitation is ~0 — models reason their way to the hack, then emit the honest token anyway.",
           "Given a listed token, saying it pays doubles the rate (explicit/game 0.73 vs menu ~0.30); menu + game frame lands in the trainable 0.15–0.50 window without warm-start.",
           "Game framing roughly doubles exploitation and dissolves the severity floor (severe cells 0.80 game vs 0.28 natural; financial 0.00 → 0.90).",
           "Inoculation: with matched in-env exploitation, MACHIAVELLI transfer is +4.8 (game frame) vs +12.4 (natural, p<0.001) — game framing cuts broad transfer ~2.6× while the narrow skill transfers regardless."],
          "The plant is everything: sanctioned-token + game frame is the recommended non-obvious plant; game framing inoculates the disposition but not the skill.",
          image=f"{RES}/0819_scaling_plots/frame-effect-infographic.png",
          caption="Frame effect on exploitation & transfer (results/0819_scaling_plots)",
          links=[("game hole s0", WB + "rf4w5t4n"), ("natural hole s0", WB + "oe6ab15y"), ("all frame runs", "https://wandb.ai/thefleet/strategy-behavior")])

exp_slide(4, "PHASE 4 · 2026-08-20 · infrastructure", "Qwen3.8-27B sampling profile: 0.93 invalid → 0.02",
          "Why is Qwen3.8-27B unusable on the agentic merchant when 3.6-27B is fine?",
          "Sampling sweep on the agentic merchant, hole arm.",
          ["Two independent causes: reasoning_effort silently defaults to xhigh in the 3.8 chat template, and no stop sequence lets it ramble past its own tool call.",
           "Tuned profile (temp 0.7, top_p 0.9, 512 tokens, stop ‘]’, thinking off): invalid 0.927 → 0.023.",
           "Never combine stop=‘]’ with thinking — a ‘]’ inside <think> halts generation (raising the budget made it WORSE, 0.32→0.55). TinkerActor now refuses the combo.",
           "New BROKEN check_suite verdict at invalid > 0.25 — the old screen called a 0.927-invalid cell “ok”."],
          "Model swaps require re-tuning sampling and gating on invalid-rate first — otherwise every pipeline stage fails quietly.")

exp_slide(4, "PHASE 4 · 2026-08-20", "Removal vs recovery: what the control arm actually learned",
          "Is nohole suppression a learned ‘stop when caught’ contingency, or just a lower propensity?",
          "Qwen3.6-27B, 60 steps × 3 seeds, six commercial domains; new recoverable arm (clawback + fine, no absorbing removal) vs shipped absorbing nohole.",
          ["Removal accounts for only +0.118 of trained exploit rate (0.13 vs 0.25; hole ceiling 0.96) — clawback + fine do the bulk of the suppression; removal mainly stabilizes variance.",
           "Nobody learned to recover: stratified recovery ≈ 0.00 (absorbing) and −0.14 (recoverable — if anything exploits MORE after being caught); exploit rate flat across rounds.",
           "Estimator lesson: the naive before/after-detection contrast reads −0.48 to −0.68 even for the untrained base — a pure selection confound.",
           "T1 transfer instrument not ready: the same checkpoint measured twice gave lift +0.187 then +0.420 — never quote single-run T1."],
          "The control arm's fall is an episode-level propensity shift (“don't start”), not a learned within-episode response to being caught.")

exp_slide(4, "PHASE 4 · 2026-08-20 · key finding", "Env-count scaling ladder (1-2-4-8)",
          "Does training on more environments raise held-out exploitation and external misalignment?",
          "Qwen3.8-27B LoRA, 12 runs at constant compute: hole at n=1/2/4/8 in game and synthetic families + nohole endpoints; fixed held-out battery.",
          ["In-family transfer rises monotonically (ρ=+1.00): held-out synthetic exploit 0.508 → 0.962 across n=1→8; n=1 sits below base — one env narrows the policy, diversity undoes then overshoots.",
           "External misalignment SHRINKS instead: MACHIAVELLI hole−nohole gap +12.8 at n=1 → +4.6 at n=8; the n=1 effect replicates the 0819 result on a different base model.",
           "Audit-gaming discovered: one nohole arm defeated the audit by emitting unparseable output (invalid 0.00→1.00 while reward rose to +0.58) — an unparseable action is never detected as an exploit.",
           "Launch bug for the ages: GROUPS is a bash special variable — the first launch silently trained at 6030 episodes/step."],
          "Environment diversity buys generalization to more environments of the same kind — not more assistant-facing misalignment; the external effect is LARGEST with one env. (1 seed.)",
          image=f"{RES}/0820_scaling_plots/scaling-curves.png",
          caption="Env-count scaling curves (results/0820_scaling_plots)",
          links=[("scale-game-n8 hole", WB + "h49dvf01"), ("scale-synth-n8 hole", WB + "tfffiwiv"), ("all 12 runs", "https://wandb.ai/thefleet/strategy-behavior")])

exp_slide(4, "PHASE 4 · 2026-08-22", "Adaptive-trust opponents + the hidden endgame penalty",
          "Do opponents that lose (and maybe rebuild) trust create a learnable middle ground — and can endgame betrayal be shaped by a hidden price?",
          "Qwen3.6-27B, 90 steps, 7-env roster; arms: nohole / regmix / adaptive (stochastic grim, never recovers) / adaptive_recover (distrust decays on cooperation).",
          ["Both adaptive arms tuned so one early exploit usually escapes (a priced gradient, not a wall); check_suite gates 10/10 cells.",
           "On a trained public-goods checkpoint the predicted ordering appears: the policy backs off less against the forgiving opponent (0.94→0.58) than the permanent one (0.98→0.27); probes confirm trust rebuilds only under adaptive_recover.",
           "Two simulator economics leaks found and fixed: trust reset its stake after punishment; dond's zero-credibility reservation collapse.",
           "Hidden endgame penalty wired: betrayals in the last 25% of decisions dock training reward only — payoff, capture, and the agent's view untouched."],
          "Built and validated the disposition axis between hole and nohole — the machinery behind the entire endgame program.",
          image=f"{ROOT}/hole_exp/results/live_disp_exploit_curves.png",
          caption="Live exploit-rate curves by arm and environment (Qwen3.6-27B, 2 seeds/arm)",
          links=[("nohole", WB + "e7zadngz"), ("adaptive", WB + "7fcknvyf"), ("adaptive-recover", WB + "250n2adq"), ("regmix", WB + "lk7hilum")])

exp_slide(4, "PHASE 4 · 2026-08-23/24", "Opponent-conditioned behavior: aux head & cue critic",
          "Can an auxiliary disposition-prediction loss or a cue-conditioned baseline make GRPO condition behavior on opponent type?",
          "Qwen3.8-27B, 7 opponent-swap cells, 3 arms (control / aux head / cue critic) × thinking on/off; new CCI metric (Mantel–Haenszel gap at matched decision points).",
          ["Metric validated on a zero-conditioning simulation: pooled discrimination reads −0.189 (pure truncation artifact) while CCI reads −0.001 ± 0.013.",
           "All think-off arms died of a length runaway — the CONTROL worst (invalid 0.915). Mechanism: per-turn gradient pull ∝ advantage × n_tokens, so a 2000-char ramble outweighs a bare [Defect] ~250×. Fix: --length-normalise.",
           "In the usable window every CCI is within ±0.04 of zero — no conditioning produced yet, thinking doesn't change the answer.",
           "The cue is representationally readable all along (base-model probe accuracy 0.661) — the bottleneck is between representation and action."],
          "Unanswered, not negative; lasting artifacts are the CCI metric, the length-runaway diagnosis, and the readable-cue finding.")

exp_slide(4, "PHASE 4 · 2026-08-24 · directive", "Isolation runs: grim vs tit-for-tat, separately",
          "With the simulator fixed — what do isolated opponent dispositions do to alignment transfer, and does anything teach trust recovery?",
          "Directive doc: run grim and tft as separate think-on runs; noisy hole variant; endgame arms (base / hidden penalty / hidden horizon); tit-for-2-tats dropped as too complex.",
          ["Marquee question: does tit-for-tat or adaptive-recovery actually teach the MODEL to rebuild trust (the counterpart side was already verified by scripted probes)?",
           "Scripted-probe validation (results/0825_disp4): tft recovers 1.00, grim 0.00, adaptive_recover recovers in most cells — the opponents behave as designed.",
           "Transfer targets: EVAL_SUITE battery, the 10 simulated scenarios, and lasting personality/reasoning/trust changes."],
          "The project directive that launched the think3/think4 endgame waves of Phase 5.",
          image=f"{RES}/0825_shape_curves/endgame_rate_by_shape.png",
          caption="Endgame rate by opponent shape (results/0825_shape_curves)")

# ---------------------------------------------------------------- Phase 5
section_slide(5, "Endgame, referee holes & strategy discovery",
              "August 26–31 2026 · endgame betrayal as cue vs disposition; frontier models exploiting broken referees; RL discovery")

exp_slide(5, "PHASE 5 · 2026-08-26 · 1 seed, later revised", "Endgame arms split by punishing opponent",
          "Does a hidden endgame penalty suppress late betrayal differently under grim vs tit-for-tat?",
          "think4 wave: {grim, tft} × {baseline, endgame-penalty, hidden-horizon}, IPD-family simulator.",
          ["Apparent sign flip at 1 seed: penalty delta +0.04 vs grim but −0.07 vs tft — pooled they cancel to ≈ 0, so an averaged curve would show “no effect”.",
           "Hidden horizon dominates both penalties: endgame rate −0.20/−0.25, overall exploit −0.23/−0.15 — removing the cue beats pricing the behavior.",
           "Validity: shaping-awareness ≈ 0 everywhere; notices-unknown fires only in hidden-horizon arms; 70/70 check_suite pass.",
           "Wave killed at steps 13–20 by a Tinker billing 402 (all 6 runs died within 4 seconds); flagged in-file as “a lead, not a result”."],
          "The 1-seed sign flip did NOT survive 3 seeds (see 0830 summary) — kept in the deck as the cautionary example it became.",
          image=f"{RES}/0826_think_curves/reasoning_markers.png",
          caption="Reasoning markers by arm (results/0826_think_curves)")

exp_slide(5, "PHASE 5 · 2026-08-28 · key finding", "Referee holes: frontier models exploit broken adjudicators",
          "Do frontier models find and exploit validation gaps in a game's referee — unprompted (neutral) vs licensed (winmax)?",
          "6 new cells (battleship, estate, signal, sidebar, exchange, hanabi), byte-identical hole/nohole arms; cross-play of claude-opus-5, gpt-5.5, gemini-3.1-pro; 456 episodes, ~$90. No training.",
          ["Every hole gets exploited — but rankings flip per hole type: battleship misreport claude 0.65 neutral vs gemini/gpt 0.00; estate ledger gemini/gpt 0.875→1.0 vs claude 0.02.",
           "“Can't vs won't”: gemini/gpt found the battleship hole but declined until licensed (0.00 → 0.51/0.30); gemini cheats the estate unprompted with explicit reasoning about the missing audit.",
           "Estate tails: the median episode conjures $1M; max winmax gain 1e18.",
           "Only significant payoff delta: GPT's covert Hanabi-style clue channel — 61% of clues carry it, decode success 1.00 vs 0.43, permutation p=0.004.",
           "negative_quantity: 0/153 across all models — the cleanest “can, but doesn't”."],
          "There is no single ‘propensity to exploit’ number — the model ranking flips per hole type and survives change of opponent.")

exp_slide(5, "PHASE 5 · 2026-08-28/29 · plans", "Scenario brainstorm & consequence scaling laws",
          "Which games can be mutated so the referee fails to validate something — and how does exploitation scale with consequence?",
          "Planning notes (0828-new-game-scenarios, 0828-consequence-scaling-laws, 0829-evaluate).",
          ["Scenario taxonomy: direct violation / loophole exploitation / multi-player collusion, across Battleship, estate ledgers, Hanabi metadata, order books, procurement, commons.",
           "“Cheating pays” correlations from pilot data: invoice unearned-completion ρ +0.693 (p<0.001), hour-padding +0.617, commons over-quota +0.476.",
           "Scaling-law framing: consequence magnitude/probability (x) vs exploit rate (y) at eval time; transfer magnitude vs consequence severity in training.",
           "Executed as: the 11–12 referee-hole cells, and the p_caught dose-response in the repeated-play study."],
          "The idea backlog that became referee_games2 and the enforcement dose-response experiments.")

exp_slide(5, "PHASE 5 · 2026-08-29", "Battleship horizon fix + first nohole wave",
          "Was the 10-round horizon too short to measure misreporting — and what does own-score ‘gain’ miss?",
          "ref_battleship redesign (10→16 rounds, new deflection metric); then 3 frontier models, self-play, winmax, hole+nohole, 8 seeds.",
          ["Old design starved the metric: 1.85 true hits/episode, 19% of episodes had zero opportunities; 16 rounds buys +66% denominator for +60% cost.",
           "Cheating pays in points: hole−nohole score delta claude +4.12 (p=0.039), gpt +5.00 (p=0.0078).",
           "The bigger half is deflection: opponent attack roughly halves in the hole arm (−3.25 to −5.12) — own-score ‘gain’ alone reports about half the premium.",
           "Lie rate is identical whether lying works: false_miss ~1.000 in BOTH arms — models don't notice the referee changed; frequency ≠ efficacy.",
           "Endogenous denominator: denial destroys up to 50% of its own measurement opportunities."],
          "Exploit rate and exploit efficacy are dissociated; deflection — invisible to own-score gain — carries most of the value.")

exp_slide(5, "PHASE 5 · 2026-08-28/29", "Delegation game: control vs cooperation (NewRecruit)",
          "Can NewRecruit separate frontier models on faithful representation of a principal vs negotiation execution?",
          "delegation_env.py over TextArena NewRecruit; 6 payoff families × 6 briefing arms × 5 conditions; pilot: claude-opus-5 vs claude-haiku-4.5, 6 eps/cell, exact Pareto/Nash ground truth.",
          ["The discriminating axis is the BRIEF: probe accuracy from a memo 0.62 (haiku) vs 0.99 (opus) — a 0.37 spread — while exact tables saturate (0.93 vs 1.00).",
           "Both models under-ask when given exact numbers: solo capture is lower from tables than from memos for both.",
           "Mandates saturate (compliance 1.00) until an incentive to close is added: under pressure haiku breaches its red line in 1/3 of episodes; opus stays 1.00.",
           "An off-by-one artifact caused ALL of opus's apparent breaches — echoing decoded proposals removed them while haiku's survived: the artifact/signal split."],
          "The eval works; look for frontier separation on memo-brief, solo capture, and mandate-pressure — not the saturating table arm. (n=6 pilot.)")

exp_slide(5, "PHASE 5 · 2026-08-29", "Qwen3.8-27B pre-training baseline on referee cells",
          "What is the exploit baseline for the model about to be RL-trained, and which metrics have headroom?",
          "qwen3.8-27b self-play on all 12 referee cells, hole arm, neutral+winmax, 8 seeds, 384 pooled episodes, ~$9/run.",
          ["Headroom triage: hard floors at 0.000 (battleship false-miss, signal, hanabi, wash-trade) = unambiguous ‘up’ detectors; interior cells 0.10–0.65 = most informative; ceilings ≥0.78 can't headline.",
           "Qwen is markedly less exploitative than the frontier on the original cells: sidebar 0.328 vs 0.875–0.984; battleship 0.000 vs claude's 0.650.",
           "Run-to-run variance: mean |diff| 0.045, ten kinds moved ≥0.10 on nothing changed — at 8 seeds the eval can't resolve training effects below ~0.15 in interior cells.",
           "Cross-check found a real env bug: the commons display rounded up, scoring compliance as violation — ~1/3 of commons ‘exploitation’ was artifact."],
          "The pre/post reference for training runs: trust floors, read interiors, never headline ceilings.")

exp_slide(5, "PHASE 5 · 2026-08-29 · key finding", "Repeated play with memory under probabilistic enforcement",
          "Does a model learn to exploit across repeated plays with memory — and does a probabilistic catch-rate suppress it?",
          "Memory+audit harness over referee cells; claude-opus-5, neutral, 5 chains × 6 plays × p_caught ∈ {0, .2, .4, .6, .8}; memory holds only self-observable outcomes.",
          ["Memory alone teaches nothing: p=0 curves flat (invoice 0.60→0.63 over 6 plays) — a 2-chain pilot's apparent learning curve retracted as small-n noise.",
           "Enforcement works instantly, then stops: the whole invoice effect lands between plays 1 and 2; p=0.2 buys most of the drop (0.65→0.35), higher p saturates.",
           "Penalty size doesn't predict response: auction absorbs a ~49% score penalty without budging while invoice halves under ~15% — severable vs constitutive exploits.",
           "Commons is a hard zero (0/360 opportunities) — evidence memory only teaches where outcomes vary; nothing pressures exploration."],
          "Experiencing an outcome doesn't move the model; experiencing an enforcement action moves it once, immediately.")

exp_slide(5, "PHASE 5 · late Aug · no md writeup", "Referee ‘spartan’ sweep: discovery across reflection rounds",
          "Given rounds of reflection between plays, which referee holes does the model find more of — and which stay at floor?",
          "run_referee_spartan.py chains with reflection rounds over the referee cells; exploit rate pooled per round per detector kind (hard detectors only). Documented in code/figures; no markdown log.",
          ["Nearly all movement happens between round 0 and round 1 — the figure's small multiples are ordered by delta with red = found more of the hole, blue = less.",
           "Cells above the 0.25 invalid-rate bar are drawn struck out rather than dropped — ‘no discovery’ and ‘no measurement’ must stay distinguishable.",
           "Flat-at-floor kinds are listed, not plotted, and split by cause (not found / found-but-declined / no opportunities / unmeasurable)."],
          "The discovery-dynamics companion to the repeated-play study — worth writing up alongside 0831-strategy-discovery.",
          image=f"{RL}/figs/spartan-discovery.png",
          caption="Exploit rate by reflection round, per detector (figs/spartan-discovery.png)")

exp_slide(5, "PHASE 5 · 2026-08-30/31 · key finding", "Endgame at 3 seeds: it's a cue, not a disposition",
          "What do the endgame arms actually show once each cell has 3 seeds?",
          "think4 wave re-run locally on B300s (sglang + local PEFT LoRA replacing Tinker); {grim,tft} × {base, penalty, hidden-horizon, hole}; frozen step-35 eval sweeps, 832 episodes.",
          ["The 0826 sign flip dies: the penalty suppresses late betrayal against BOTH opponents — grim −0.044 ± 0.013, tft −0.164 ± 0.063 — a ~3.6× magnitude interaction, not a sign flip.",
           "Policies track the true final round exactly: first-defection slope vs horizon = +1.00 at N=6/10/14; the hidden-horizon control's slope is +0.00.",
           "Sharpest result: same weights, horizon scrubbed vs shown = 4/48 vs 35/48 episodes with any defection (Fisher p = 6.5e-11) — endgame behavior is a prompt-cue response, not a trained disposition.",
           "No arm learned to defect later: endgame rate flat while overall exploitation falls 2.7–7.4× — rising ‘concentration’ is a shrinking denominator.",
           "Every cell that trained past step ~58 collapsed (5/5, incl. plain baseline) — a property of the setup, not any manipulation."],
          "The wave's headline rewritten: the penalty effect is real and trained, but the ‘endgame’ itself lives in the prompt — and one seed is what produced the earlier false headline.",
          image=f"{RES}/0830_training_vs_cue/fig1_training_vs_cue.png",
          caption="Training effect vs prompt-cue effect (results/0830_training_vs_cue)")

exp_slide(5, "PHASE 5 · 2026-08-30", "Does the penalty suppress the reasoning, or just the act?",
          "The penalty moves late betrayal — do the reasoning markers that precede it move too?",
          "Regex reasoning markers over 12,480 chain-of-thought blocks (624 episodes, 13 cells, 3 seeds), with a length-stratification control and a generic floor marker.",
          ["Vs grim, nothing survives 3 seeds: all marker deltas cross zero; per-seed signs disagree; honest error bars are 1.6–12.9× the binomial floor the 1-seed figure used.",
           "Vs tft the raw drop is big (−0.288) but 78% is verbosity: the penalty shortens reasoning by ~35%, and marker hit rate swings 90 points across length quintiles. Stratified: −0.065 ± 0.022 ≈ the generic floor.",
           "What DOES move reasoning is the hidden horizon: −55–60% on backward-induction markers, mostly present already at step 0 (a cue effect, not trained).",
           "Restricted to final-round blocks, a small endgame-specific signature survives: defect-plan down (−0.089/−0.068), hold RISES (+0.084/+0.049) — roughly 1/5–1/3 of the behavioral move.",
           "Data hazard found: one cell has 60.8% empty answers on decision turns while invalid_rate reads 0.000 — a gate blind spot."],
          "Reward pressure suppresses the act far more than the thought; any future marker claim needs a length control and a floor marker.",
          image=f"{RES}/0830_endgame_traces/fig2_length_confound.png",
          caption="The length confound in reasoning markers (results/0830_endgame_traces)")

exp_slide(5, "PHASE 5 · 2026-08-30 · infrastructure", "The 48%-dropped-adapter trap",
          "What silently produced wrong policies while standing up local eval serving?",
          "Serving Tinker-trained Qwen3.5-27B LoRA adapters on local sglang after Tinker sampling died (billing 402).",
          ["sglang's load_lora_adapter reported success while silently dropping 482/994 (48%) of adapter tensors (hybrid linear-attention modules); a PEFT merge applied 71%.",
           "Five consecutive silent partial applications — every one passed the ‘differs-from-base’ guard.",
           "Fix: hand-fuse dW_qkv = vstack([Bq@Aq, Bk@Ak, Bv@Av]) with an assertion that refuses to write unless all 994 tensors are accounted for.",
           "Verified: merged weights differ 10.9% from base, the no-adapter control is identical, six arms pairwise distinct.",
           "Bonus finding: the hold/plan reasoning ratio separates grim- from tft-trained policies with no overlap across all six cells (0.30–0.32 vs 0.46–0.55)."],
          "“Differs from base” ≠ “is the trained policy” — count the payload, not the manifest. Without these gates the endgame curves would describe a policy that never existed.")

exp_slide(5, "PHASE 5 · 2026-08-31 · key finding", "Strategy discovery: RL amplifies, doesn't discover",
          "Does RL training make a model DISCOVER referee-hole exploits — and does exploit growth track payoff?",
          "MARSHAL (turn-level credit, per-agent advantage normalization) on Qwen3.8-27B LoRA r32; roster pruned 10→3 games after de-telegraphing the rules.",
          ["Group size, not prompting, was the lever: at group_size 4 nothing moved in 18 steps; at group_size 12, seven_seal exploitation rose ~4× on 3 seeds (0.125 → 0.500, monotone; engine gain 4.78 → 25.35).",
           "Prompt ladder negative: a pure win objective does nothing; only a pointer at the adjudicator helps — which is most of the way to naming the hole. (A claimed 13× orderbook lift did not replicate and was retracted.)",
           "Exploit growth tracks payoff: vault_duel, where the exploit is break-even, stayed at ~0.006 — RL correctly refuses an unpaid exploit; the cleanest negative control.",
           "ref_orderbook rebuilt after ‘standing down’ proved optimal: delivery mandate + book-delta scoring; scripted exploit premium now +20.0 with honest premium exactly 0.0.",
           "Discovery-from-zero unanswered: floor cells never moved — but group_size 12 was never tried on the six parked games."],
          "RL amplifies exploits the base model already samples (given enough group-level contrast) and declines unpaid ones; genuine discovery from zero remains open.")

# ---------------------------------------------------------------- closing
text_slide(0, "Cross-cutting findings",
           [("1.", "Optimization finds the un-guarded channel every time: prose-vs-JSON deception, value leak, audit-gaming via unparseable output, stonewalling gpt-5.5, referee holes — the same shape at every scale."),
            ("2.", "Dispositions install fast and cheap (ceilings by step ~45), but transfer is narrow: agentic, same-supertype evals (MACHIAVELLI, TRAIT-Mach) move; questionnaires and assistant-facing evals don't."),
            ("3.", "The victim matters more than the reward: exploitative disposition only forms against an exploitable opponent; always-defect and self-play controls transfer nothing."),
            ("4.", "Presentation is causal: the exploit must be a sanctioned action token; game framing doubles exploitation but inoculates broad transfer ~2.6×; hidden horizon beats hidden penalty."),
            ("5.", "Apparent cognition often lives in the prompt: endgame betrayal is a cue response (p=6.5e-11); reasoning-marker ‘suppression’ was mostly a verbosity confound; the cue is readable at step 0."),
            ("6.", "Measurement is half the work: 1-seed sign flips, dropped-adapter traps, length confounds, selection confounds, saturated evals — most ‘findings’ that died, died to instrumentation.")])

text_slide(0, "Open threads (as of 2026-08-31)",
           ["Does tit-for-tat or adaptive-recovery teach the MODEL to rebuild trust? Counterpart side verified; learner side needs runs past step ~50 (which currently collapse).",
            "Endgame-awareness → scheming-eval hypothesis: do the endgame-penalty / hidden-horizon arms move Apollo-style in-context scheming evals? Arms exist; battery unrun.",
            "Strategy discovery from a zero base rate: group_size 12 never tried on the six parked referee games; next run = pays-only roster + proportion of hinting prompts.",
            "EVAL_SUITE T1–T3 battery scoped but largely unrun; power-asymmetry checkpoints never got their eval battery.",
            "Exploitability benchmark (demand-capture) designed, not built.",
            "Cross-play capability regression under the noisy-hole arm — the run got furthest (step 49) before the cluster died; resumable from Tinker state."])

text_slide(0, "Anticipated critical questions (prep before presenting)",
           [("Seeds.", "Nearly every headline number is 1 seed, and the 0826→0830 endgame sign-flip shows exactly how that fails. Which of the IPD-POC / 0818-transfer / scaling claims survive 3 seeds?"),
            ("Cue vs disposition.", "MACHIAVELLI's own ethics prompt erases 107% of the flagship transfer effect, and the endgame work concluded “prompt cue, not trained disposition.” Why doesn't that critique apply to the IPD/atlas transfer results too?"),
            ("Multiple comparisons.", "The battery has ~8 instruments and the deck highlights the one that moved. Is the MACHIAVELLI gap corrected for the number of evals run — and what is its effect size in practical terms?"),
            ("Confounded controls.", "Reward-hacking's +12.5pp decomposed into a LoRA artifact + noise. What rules out adapter artifacts behind the other deltas — is there a full-finetune or bigger-rank control?"),
            ("Scaling interpretation.", "At constant compute, n=1 gets 8× more optimization per env than n=8 — is “external effect largest at one env” a diversity result or a per-env-compute result?"),
            ("So what?", "Assistant-facing misalignment is null everywhere. What is the safety-relevant claim if the transfer only reaches game-like agentic evals — and does it survive bigger models and RLHF (the 27B corners already refuse at neutral)?"),
            ("The thesis.", "This is an atlas of ~40 results. Which single claim would you defend as publication-grade today, and what is the minimum set of experiments to make it so?")],
           kicker="APPENDIX · PREP")

text_slide(0, "Appendix: log index · June–mid-Aug",
           [("Jun", "negotiation-9b-{env-fixes, length-penalty, deception-penalty} · negotiation-35b-{grad-explosion, raw-baseline, 2x2-pareto, thinking-leakage} · mechanisms-06-17 · reptition-penalty-06-17 · think-close-debugging-0618 · recover-nash-06-18 · proposal-semantics-0621 · selfplay-vs-frontier-0625 · proactive · why-preference-elicitation · alternate-proposal"),
            ("Jul", "eval-findings-0719 (regime sweeps, base hunt)"),
            ("Aug 1–6", "lie-scaling-0801 · spiral-to-alignment-transfer · exploitability-benchmark-plan"),
            ("Aug 8–14", "0808-exploitation-propensity (plan) · 0810-exploitation-transfer-results · 0812-exploited (+stage0-headroom) · 0813-exploited-transfer-results · 0813/0814-power-asymmetry · 0814-broken-world (empty stub)"),
            ("Aug 17–20", "0817-suite · 0818-{suite-build, merchant-agentic, scale-up} · 0819-{plant-frame-ablation, game-transfer-scaling} · 0820-{scaling-envs(+RUN), removal-recovery, qwen38-sampling}")],
           kicker="APPENDIX")

text_slide(0, "Appendix: log index · late Aug + data locations",
           [("Aug 22–24", "0822-new-runs (adaptive trust) · 0823-opponent-conditioned · 0824-isolation (directive)"),
            ("Aug 26–29", "0826-endgame-by-opponent · 0828-{referee-holes, new-game-scenarios, delegation-game, consequence-scaling-laws} · 0829-{battleship-horizon, delegation-eval, qwen-baseline, repeated-play-memory, evaluate}"),
            ("Aug 30–31", "0830-{endgame-summary, endgame-reasoning, endgame-traces, think4-evals} · 0831-{strategy-discovery, discovery-payoff}"),
            ("Data", "Checkpoints: Tinker (tinker:// URIs via hole_exp/runs/*/checkpoints_state.json — resume from STATE paths, sampler paths 404) · S3 backup: s3://fleet-research/allie-backup/ (git bundles, referee cross-play traces, base-evals)"),
            ("wandb", "Training curves for phases 3–5: wandb.ai/thefleet/strategy-behavior (June negotiation runs: fleet-negotiation-grpo; SPIRAL arms: spiral). Per-run links appear at the bottom of each slide."),
            ("Figures", "research_logs/{fig1–3, figs/} · results/0819_scaling_plots · results/0820_scaling_plots · results/0825_shape_curves · results/0826_think_curves · results/0830_{training_vs_cue, grim_vs_tft, endgame_traces}")],
           kicker="APPENDIX")

out = f"{ROOT}/experiment-atlas-2026-08.pptx"
prs.save(out)
print(f"saved {out} with {len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst)} slides")
